import io

import pytest
from pptx import Presentation
from pptx.util import Inches

from backend.app.services.knowledge_service import _extract_text


def _make_pptx() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Northwind Labs"
    slide.placeholders[1].text = "Real-time data infrastructure"
    # A table whose cell text must also be extracted.
    table_slide = prs.slides.add_slide(prs.slide_layouts[5])
    tbl = table_slide.shapes.add_table(2, 2, Inches(0.5), Inches(1.5), Inches(4), Inches(1)).table
    tbl.cell(0, 0).text = "Vendor"
    tbl.cell(0, 1).text = "Latency"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_pptx_pulls_slide_and_table_text():
    text = _extract_text(_make_pptx(), "deck.pptx")
    assert "Northwind Labs" in text
    assert "Real-time data infrastructure" in text
    assert "Vendor" in text and "Latency" in text


def test_extract_plaintext_types():
    assert "hello world" in _extract_text(b"hello world", "notes.md")
    assert "a,b,c" in _extract_text(b"a,b,c", "data.csv")


def test_extract_unsupported_type_raises():
    with pytest.raises(ValueError):
        _extract_text(b"x", "image.png")
