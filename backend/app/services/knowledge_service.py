import io
import uuid
import asyncpg
from openai import AsyncOpenAI
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from backend.app.config import settings
from backend.app.services.storage_service import upload_object

openai_client = AsyncOpenAI(api_key=settings.openai_api_key)


def _storage_key(company_id: str, doc_id: uuid.UUID, filename: str) -> str:
    """Tenant-scoped object key for a raw uploaded document."""
    return f"{company_id}/{doc_id}/{filename}"

CHUNK_SIZE = 512
CHUNK_OVERLAP = 50

def _extract_pptx_text(content: bytes) -> str:
    """Pull text from every slide: titles, body text frames, table cells, and
    speaker notes. One paragraph per logical block so chunking stays coherent."""
    prs = Presentation(io.BytesIO(content))
    blocks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(
                    p.text for p in shape.text_frame.paragraphs if p.text.strip()
                )
                if text.strip():
                    blocks.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        blocks.append(" | ".join(cells))
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        if notes.strip():
            blocks.append(notes)
    return "\n\n".join(blocks)


def _extract_text(content: bytes, filename: str) -> str:
    name = filename.lower()
    if name.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(content))
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)
    elif name.endswith((".docx", ".doc")):
        doc = DocxDocument(io.BytesIO(content))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    elif name.endswith(".pptx"):
        return _extract_pptx_text(content)
    elif name.endswith((".txt", ".md", ".markdown", ".csv")):
        return content.decode("utf-8", errors="ignore")
    raise ValueError(f"Unsupported file type: {filename}")

def _chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Pack paragraphs into word-bounded chunks with a sliding overlap.

    Paragraphs longer than `chunk_size` are pre-split into windows so no single
    chunk balloons past the target — oversized chunks dilute embedding relevance
    and hurt retrieval precision.
    """
    raw_paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    stride = max(1, chunk_size - overlap)
    chunks: list[str] = []
    current_words: list[str] = []

    for para in raw_paragraphs:
        words = para.split()
        # A paragraph longer than the target is emitted as strict sliding windows
        # (each <= chunk_size), bypassing the packing buffer so the overlap carry
        # can't push a window over the cap.
        if len(words) > chunk_size:
            if current_words:
                chunks.append(" ".join(current_words))
                current_words = []
            for i in range(0, len(words), stride):
                chunks.append(" ".join(words[i : i + chunk_size]))
                if i + chunk_size >= len(words):
                    break
            continue
        if current_words and len(current_words) + len(words) > chunk_size:
            chunks.append(" ".join(current_words))
            current_words = current_words[-overlap:] if overlap else []
        current_words.extend(words)

    if current_words:
        chunks.append(" ".join(current_words))
    return chunks

async def _embed_texts(texts: list[str]) -> list[list[float]]:
    response = await openai_client.embeddings.create(
        model="text-embedding-3-small",
        input=texts,
    )
    return [item.embedding for item in response.data]

async def ingest_document(
    conn: asyncpg.Connection,
    content: bytes,
    filename: str,
    company_id: str,
    doc_type: str,
) -> dict:
    doc_id = uuid.uuid4()
    namespace = f"{company_id}:{doc_type}"

    # Persist the raw file to object storage when configured (skipped in local
    # dev where no bucket is set). The key is stored for later retrieval.
    s3_key = None
    if settings.s3_bucket_name:
        s3_key = _storage_key(company_id, doc_id, filename)
        await upload_object(content, s3_key)

    await conn.execute(
        """INSERT INTO knowledge_documents (id, company_id, filename, doc_type, status, s3_key)
           VALUES ($1, $2, $3, $4, 'pending', $5)""",
        doc_id, uuid.UUID(company_id), filename, doc_type, s3_key,
    )

    try:
        text = _extract_text(content, filename)
        chunks = _chunk_text(text)
        embeddings = await _embed_texts(chunks)

        async with conn.transaction():
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                await conn.execute(
                    """INSERT INTO document_chunks
                       (document_id, company_id, namespace, chunk_index, content, embedding)
                       VALUES ($1, $2, $3, $4, $5, $6::vector)""",
                    doc_id, uuid.UUID(company_id), namespace, i, chunk, str(embedding),
                )
            await conn.execute(
                "UPDATE knowledge_documents SET status='indexed', chunk_count=$1 WHERE id=$2",
                len(chunks), doc_id,
            )

        return {"document_id": str(doc_id), "chunks": len(chunks), "status": "indexed"}

    except Exception as exc:
        await conn.execute(
            "UPDATE knowledge_documents SET status='failed' WHERE id=$1", doc_id
        )
        raise exc
